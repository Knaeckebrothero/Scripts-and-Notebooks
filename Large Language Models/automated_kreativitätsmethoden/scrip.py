import openai
import wikipediaapi
import os
import json
import csv
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches


def api_request(method_description):
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo-16k",
        messages=[
            {
                "role": "system",
                "content": "Du bist ein hilfreicher Assistent. Du bekommst von deinem Vorgesetzen „Kreativitätsmethoden“ und sollst zu diese eine kurze und Anschauliche Beschreibung und ein Beispiel erstellen. Da die Inhalte maschinell in eine Präsentation eingefügt werden sollen, wirst du sie als JSON-String formatieren. Bedeutet das zu jeder Methode die folgenden zwei Attribute benötigt werden: description, example. Deine Antwort sollte Faktenbasiert und in deutscher Sprache (unabhängig vom gegebenen Inhalt) formuliert sein. Verwende für die Beschreibung lediglich gesicherte Informationen aus dem dir zur Verfügung gestellten Material. Bedeutet, du sollst unter keinen Umständen dir Inhalte ausdenken, oder halluzinieren."
            },
            {
                "role": "user",
                "content": method_description
            }],
        temperature=1,
        max_tokens=7000,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
    )
    
    try:
        # Save in json file
        with open('response.json', 'w') as outfile:
            json.dump(response['choices'][0]['message']['content'], outfile)

        # Return the content
        return json.loads(response['choices'][0]['message']['content'])
    except:
            try:
                new_response = openai.ChatCompletion.create(
                    model="gpt-3.5-turbo-16k",
                    messages=[
                        {
                            "role": "system",
                            "content": "Du bist ein hilfreicher Assistent. Du bekommst von deinem Vorgesetzen „Kreativitätsmethoden“ und sollst zu diese eine kurze und Anschauliche Beschreibung und ein Beispiel erstellen. Da die Inhalte maschinell in eine Präsentation eingefügt werden sollen, wirst du sie als JSON-String formatieren. Bedeutet das zu jeder Methode die folgenden zwei Attribute benötigt werden: description, example. Deine Antwort sollte Faktenbasiert und in deutscher Sprache (unabhängig vom gegebenen Inhalt) formuliert sein. Verwende für die Beschreibung lediglich gesicherte Informationen aus dem dir zur Verfügung gestellten Material. Bedeutet, du sollst unter keinen Umständen dir Inhalte ausdenken, oder halluzinieren."
                        },
                        {
                            "role": "user",
                            "content": "Leider hat der JSON String einen fehler bitte überprüfe ihn und korregiere den Fehler! Damit die JSON Datei richtig formatiert ist, muss sie folgende Struktur haben: [{\"description\": \"Beschreibung\", \"example\": \"Beispiel\"}. Der Fehlerhafte JSON STRING: " + response['choices'][0]['message']['content']
                        }],
                    temperature=1,
                    max_tokens=7000,
                    top_p=1,
                    frequency_penalty=0,
                    presence_penalty=0
                    )
                # Save in json file
                with open('response.json', 'w') as outfile:
                    json.dump(new_response['choices'][0]['message']['content'], outfile)

                # Return the content
                return json.loads(new_response['choices'][0]['message']['content'])
            except:
                return "Error"

def get_wikipedia_article_content(german_link, english_link, user_agent):
    # Create a Wikipedia object with a specific user agent
    wiki_wiki_de = wikipediaapi.Wikipedia(user_agent=user_agent, language='de')
    wiki_wiki_en = wikipediaapi.Wikipedia(user_agent=user_agent, language='en')

    # Extract the title from the German link
    german_title = german_link.split('/')[-1]
    page_de = wiki_wiki_de.page(german_title)

    # Check if the German page exists
    if page_de.exists():
        return page_de.text

    # If German page doesn't exist, try the English link
    if english_link:
        english_title = english_link.split('/')[-1]
        page_en = wiki_wiki_en.page(english_title)
        if page_en.exists():
            return page_en.text

    return "Page not found"

# Set the OpenAI API key
load_dotenv()
openai.api_key = os.getenv("OPENAIKEY")
user = os.getenv("USERAGENT")

with open('Kreativitätsmethoden.csv') as csv_file:
    reader = csv.DictReader(csv_file)
    methods = [row for row in reader]

    presentation = Presentation()

    for method in methods:
        # Counter
        print(method)

        # Get the content from Wikipedia
        german_link = method['German Wikipedia Link']
        english_link = method['English Wikipedia Link']
        content = get_wikipedia_article_content(german_link, english_link, user)
        response = api_request(content)
        
        if response == "Error":
            continue

        # Get the description and example from the response
        description = response['description']
        example = response['example']

        # Create a new slide with a title and two content boxes
        slide_layout = presentation.slide_layouts[1] # Use layout 1 for title and content
        slide = presentation.slides.add_slide(slide_layout)

        # Set the title
        title = slide.shapes.title
        title.text = method['Method']

        # Add description on the left
        description_box = slide.placeholders[1]
        description_box.text = "Description:\n" + description

        # Add example on the right
        example_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(5), Inches(4))
        example_frame = example_box.text_frame
        example_frame.text = "Example:\n" + example

        # Save the presentation
        presentation.save('Kreativitätsmethoden.pptx')
